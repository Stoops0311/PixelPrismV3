"""
PixelPrism — Enterprise Customer Kickoff Deck.

Studio Brutalist DS-2 design system rendered to PPTX via python-pptx.
Run:  python3 build_deck.py
Output: PixelPrism-Enterprise-Kickoff.pptx (next to this script).
"""

from __future__ import annotations

import os
from dataclasses import dataclass

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


# ============================================================================
# THEME
# ============================================================================

class Color:
    BACKGROUND = RGBColor(0x07, 0x1a, 0x26)
    CARD = RGBColor(0x0e, 0x28, 0x38)
    MUTED_BG = RGBColor(0x0b, 0x22, 0x30)
    POPOVER = RGBColor(0x16, 0x33, 0x44)
    GOLD = RGBColor(0xf4, 0xb9, 0x64)
    CORAL = RGBColor(0xe8, 0x95, 0x6a)
    LIME = RGBColor(0xa4, 0xf4, 0x64)
    CYAN = RGBColor(0x64, 0xdc, 0xf4)
    LIGHT_GOLD = RGBColor(0xf4, 0xd4, 0x94)
    MUTED = RGBColor(0x6d, 0x8d, 0x9f)
    FOREGROUND = RGBColor(0xea, 0xee, 0xf1)
    SECONDARY_FG = RGBColor(0xd4, 0xdc, 0xe2)
    BLACK = RGBColor(0x00, 0x00, 0x00)


def _blend(fg: RGBColor, bg: RGBColor, alpha: float) -> RGBColor:
    """Composite fg onto bg at given alpha (0–1). PPTX has no alpha on solid fills."""
    r = round(bg[0] + (fg[0] - bg[0]) * alpha)
    g = round(bg[1] + (fg[1] - bg[1]) * alpha)
    b = round(bg[2] + (fg[2] - bg[2]) * alpha)
    return RGBColor(max(0, min(255, r)), max(0, min(255, g)), max(0, min(255, b)))


class GoldOn:
    BG_06 = _blend(Color.GOLD, Color.BACKGROUND, 0.06)
    BG_08 = _blend(Color.GOLD, Color.BACKGROUND, 0.08)
    BG_12 = _blend(Color.GOLD, Color.BACKGROUND, 0.12)
    BG_22 = _blend(Color.GOLD, Color.BACKGROUND, 0.22)
    CARD_12 = _blend(Color.GOLD, Color.CARD, 0.12)
    CARD_22 = _blend(Color.GOLD, Color.CARD, 0.22)


class Font:
    DISPLAY = "Neue Montreal"
    HEAD = "Neue Montreal"
    BODY = "General Sans"
    MONO = "JetBrains Mono"


# Slide canvas
SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN_L = 0.6
MARGIN_R = 0.6


# ============================================================================
# LOW-LEVEL HELPERS
# ============================================================================

def _kill_shadow(shape):
    """Override theme shadow on a shape by injecting an empty <a:effectLst/>."""
    spPr = shape._element.find(qn("p:spPr"))
    if spPr is None:
        return
    for el in spPr.findall(qn("a:effectLst")):
        spPr.remove(el)
    etree.SubElement(spPr, qn("a:effectLst"))


def _apply_letter_spacing(run, em_fraction: float):
    """Set letter spacing as a fraction of em (font size). 0.10 == 0.10em."""
    if not em_fraction:
        return
    size_pt = run.font.size.pt if run.font.size else 14
    spc_hundredths = int(em_fraction * size_pt * 100)
    rPr = run._r.find(qn("a:rPr"))
    if rPr is None:
        rPr = etree.SubElement(run._r, qn("a:rPr"))
        run._r.insert(0, rPr)
    rPr.set("spc", str(spc_hundredths))


def set_slide_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w_pt=0.75):
    """Add a sharp-edged rectangle. fill/line accept RGBColor or None (transparent)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    _kill_shadow(shape)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w_pt)
    return shape


def add_text(
    slide,
    x, y, w, h,
    text: str,
    font: str = Font.BODY,
    size: float = 14,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor | None = None,
    align: str = "left",
    anchor: str = "top",
    letter_spacing: float = 0.0,
):
    """Add a textbox with a single run. Returns the textbox shape."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    if anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.alignment = {
        "left": PP_ALIGN.LEFT,
        "right": PP_ALIGN.RIGHT,
        "center": PP_ALIGN.CENTER,
    }[align]
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color
    if letter_spacing:
        _apply_letter_spacing(run, letter_spacing)
    return tb


def add_paragraphs(
    slide,
    x, y, w, h,
    lines: list[str],
    font: str = Font.BODY,
    size: float = 14,
    bold: bool = False,
    color: RGBColor | None = None,
    line_spacing: float = 1.4,
    align: str = "left",
):
    """Add a textbox containing multiple paragraphs (one per line)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    align_enum = {
        "left": PP_ALIGN.LEFT,
        "right": PP_ALIGN.RIGHT,
        "center": PP_ALIGN.CENTER,
    }[align]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align_enum
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color
    return tb


def add_bulleted_list(
    slide,
    x, y, w, h,
    bullets: list[str],
    font: str = Font.BODY,
    size: float = 14,
    color: RGBColor | None = None,
    bullet_color: RGBColor | None = None,
    line_spacing: float = 1.4,
    bullet_gap: float = 0.32,
):
    """
    Add a list where each item is a tiny gold square + paragraph of body text.
    We render each bullet as: a small gold rect + a sibling textbox.
    """
    if color is None:
        color = Color.SECONDARY_FG
    if bullet_color is None:
        bullet_color = Color.GOLD
    line_h = bullet_gap  # vertical step per item, inches
    for i, b in enumerate(bullets):
        row_y = y + i * line_h
        # Square marker
        marker_size = 0.09
        marker_y = row_y + 0.13  # nudge down to align with text baseline
        add_rect(slide, x, marker_y, marker_size, marker_size, fill=bullet_color)
        # Text
        add_text(
            slide,
            x + 0.28,
            row_y,
            w - 0.28,
            line_h * 1.4,
            b,
            font=font,
            size=size,
            color=color,
        )


def add_overline(slide, x, y, w, text: str, color: RGBColor | None = None):
    if color is None:
        color = Color.GOLD
    add_text(
        slide, x, y, w, 0.3, text.upper(),
        font=Font.HEAD, size=10, bold=True, color=color,
        letter_spacing=0.12,
    )


def add_gold_accent_bar(slide, x, y, height: float = 0.5):
    """Vertical gold accent bar — used to the left of headlines."""
    add_rect(slide, x, y, 0.045, height, fill=Color.GOLD)


def add_gold_rule(slide, x, y, length: float):
    """Horizontal gold rule, full saturation."""
    add_rect(slide, x, y, length, 0.035, fill=Color.GOLD)


def add_card(slide, x, y, w, h, fill=None, border=None):
    """A dark card with a gold-tinted hairline border, plus a dual-shadow effect."""
    if fill is None:
        fill = Color.CARD
    if border is None:
        border = GoldOn.BG_12
    # Drop-shadow layer (broad ambient): a dark rectangle offset down/right
    add_rect(slide, x + 0.04, y + 0.07, w, h, fill=RGBColor(0x03, 0x10, 0x18))
    # Contact-shadow layer (tight): tighter, slightly different shade
    add_rect(slide, x + 0.015, y + 0.025, w, h, fill=RGBColor(0x05, 0x14, 0x1d))
    # Card surface
    return add_rect(slide, x, y, w, h, fill=fill, line=border, line_w_pt=0.75)


# ============================================================================
# FRAME (recurring elements on every slide)
# ============================================================================

TOTAL_PAGES = 20


def slide_frame(prs, page_num: int):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, Color.BACKGROUND)

    # Top-right: mono page counter
    add_text(
        slide, SLIDE_W - 1.8, 0.32, 1.4, 0.3,
        f"{page_num:02d} / {TOTAL_PAGES:02d}",
        font=Font.MONO, size=9, color=Color.MUTED, align="right",
        letter_spacing=0.02,
    )

    # Bottom-left: wordmark
    add_text(
        slide, MARGIN_L, SLIDE_H - 0.45, 3, 0.3,
        "PIXELPRISM", font=Font.HEAD, size=8.5, bold=True,
        color=Color.MUTED, letter_spacing=0.20,
    )

    # Bottom-right: gold square (live-dot motif, static)
    add_rect(slide, SLIDE_W - 0.66, SLIDE_H - 0.36, 0.08, 0.08, fill=Color.GOLD)

    return slide


# ============================================================================
# SLIDE BUILDERS
# ============================================================================

def build_cover(prs, page_num):
    s = slide_frame(prs, page_num)

    # Top-left overline
    add_overline(s, MARGIN_L, 0.55, 8, "Kickoff / Enterprise Introduction")

    # Display headline
    add_text(
        s, MARGIN_L, 2.0, 12, 1.6,
        "PIXELPRISM.",
        font=Font.DISPLAY, size=92, bold=True, color=Color.FOREGROUND,
        letter_spacing=-0.02,
    )

    # Gold rule
    add_gold_rule(s, MARGIN_L, 3.55, 1.6)

    # Subtitle line 1
    add_text(
        s, MARGIN_L, 3.85, 12, 0.9,
        "Social media marketing,",
        font=Font.DISPLAY, size=38, bold=True, color=Color.SECONDARY_FG,
    )
    # Subtitle line 2
    add_text(
        s, MARGIN_L, 4.55, 12, 0.9,
        "engineered for enterprise brands.",
        font=Font.DISPLAY, size=38, bold=True, color=Color.GOLD,
    )

    # Tagline
    add_text(
        s, MARGIN_L, 5.85, 8, 0.4,
        "Sharp but Alive.",
        font=Font.BODY, size=14, italic=True, color=Color.MUTED,
    )

    # Top-right meta block
    add_text(
        s, SLIDE_W - 3.6, 0.85, 3.0, 0.3,
        "ENTERPRISE / v1.0",
        font=Font.MONO, size=9, color=Color.MUTED, align="right",
        letter_spacing=0.05,
    )


def build_problem(prs, page_num):
    s = slide_frame(prs, page_num)
    add_overline(s, MARGIN_L, 0.7, 8, "01 / The Problem")
    add_gold_accent_bar(s, MARGIN_L, 1.4, height=0.65)
    add_text(
        s, MARGIN_L + 0.25, 1.35, 11.5, 1.0,
        "Brand-consistent social at scale is broken.",
        font=Font.HEAD, size=42, bold=True, color=Color.FOREGROUND,
        letter_spacing=-0.01,
    )

    add_paragraphs(
        s, MARGIN_L, 2.9, 9.5, 2.8,
        [
            "Enterprise marketing teams juggle dozens of brands, hundreds of products, and a daily content treadmill.",
            "The existing toolchain — scattered Canva files, spreadsheet calendars, and AI tools that don't know your brand — produces inconsistent work, slow turnarounds, and assets nobody can find again next quarter.",
        ],
        font=Font.BODY, size=16, color=Color.SECONDARY_FG, line_spacing=1.5,
    )

    # Right callout card with a pull-quote
    add_card(s, 9.6, 1.5, 3.0, 4.2)
    add_text(
        s, 9.85, 1.85, 2.6, 0.35,
        "WHY NOW",
        font=Font.HEAD, size=10, bold=True, color=Color.CORAL,
        letter_spacing=0.15,
    )
    add_paragraphs(
        s, 9.85, 2.3, 2.6, 3.2,
        [
            "We built PixelPrism",
            "for the teams that",
            "have outgrown",
            "“good enough.”",
        ],
        font=Font.HEAD, size=22, bold=True, color=Color.GOLD, line_spacing=1.15,
    )


def build_what_is_pp(prs, page_num):
    s = slide_frame(prs, page_num)
    add_overline(s, MARGIN_L, 0.7, 8, "02 / Solution")
    add_gold_accent_bar(s, MARGIN_L, 1.4, height=0.65)
    add_text(
        s, MARGIN_L + 0.25, 1.35, 12, 1.0,
        "PixelPrism is a marketing OS for enterprise brands.",
        font=Font.HEAD, size=36, bold=True, color=Color.FOREGROUND,
        letter_spacing=-0.01,
    )

    add_paragraphs(
        s, MARGIN_L, 2.65, 11.5, 1.2,
        [
            "One workspace. Many brands. Every asset, every product, every post — produced, scheduled, and measured in one place.",
        ],
        font=Font.BODY, size=16, color=Color.SECONDARY_FG, line_spacing=1.5,
    )

    # Three pillars
    pillar_y = 4.0
    pillar_h = 2.4
    gap = 0.25
    pillar_w = (SLIDE_W - 2 * MARGIN_L - 2 * gap) / 3
    pillars = [
        ("I.", "SET UP YOUR BRAND", "Configure once.\nReuse forever."),
        ("II.", "MAKE THE WORK", "Generate. Refine.\nApprove."),
        ("III.", "PUBLISH & LEARN", "Ship the work.\nMeasure the lift."),
    ]
    for i, (numeral, title, desc) in enumerate(pillars):
        x = MARGIN_L + i * (pillar_w + gap)
        add_card(s, x, pillar_y, pillar_w, pillar_h)
        add_text(
            s, x + 0.25, pillar_y + 0.22, pillar_w - 0.4, 0.7,
            numeral, font=Font.DISPLAY, size=48, bold=True, color=Color.GOLD,
        )
        add_text(
            s, x + 0.25, pillar_y + 1.1, pillar_w - 0.4, 0.35,
            title, font=Font.HEAD, size=12, bold=True, color=Color.FOREGROUND,
            letter_spacing=0.12,
        )
        add_paragraphs(
            s, x + 0.25, pillar_y + 1.55, pillar_w - 0.4, 0.8,
            desc.split("\n"),
            font=Font.BODY, size=13, color=Color.MUTED, line_spacing=1.3,
        )


def build_act_divider(prs, page_num, roman: str, title: str, promise: str):
    s = slide_frame(prs, page_num)
    # Overline
    add_overline(s, MARGIN_L, 0.7, 8, f"ACT {roman} / DIVIDER")

    # Giant numeral
    add_text(
        s, MARGIN_L, 1.4, 6.5, 5.4,
        roman,
        font=Font.DISPLAY, size=320, bold=True, color=Color.GOLD,
        letter_spacing=-0.05,
    )

    # Coral accent stripe to the right of the numeral
    add_rect(s, 5.8, 2.3, 0.08, 3.4, fill=Color.CORAL)

    # Act overline + title
    add_overline(s, 6.2, 2.6, 7, f"ACT {roman}")
    add_text(
        s, 6.2, 3.0, 7, 1.0,
        title,
        font=Font.HEAD, size=44, bold=True, color=Color.FOREGROUND,
        letter_spacing=-0.01,
    )
    # Promise
    add_text(
        s, 6.2, 4.6, 7, 0.6,
        promise,
        font=Font.BODY, size=18, italic=True, color=Color.MUTED,
    )


def build_feature(
    prs,
    page_num,
    overline: str,
    title: str,
    promise: str,
    bullets: list[str],
    visual_kind: str,
    visual_payload: dict | None = None,
):
    """
    Standard feature slide.
      Left col: overline, title, promise sentence, bulleted list
      Right col: a card with one of several visual treatments
    """
    s = slide_frame(prs, page_num)
    add_overline(s, MARGIN_L, 0.7, 8, overline)
    add_gold_accent_bar(s, MARGIN_L, 1.4, height=0.65)
    add_text(
        s, MARGIN_L + 0.25, 1.35, 7.5, 1.0,
        title,
        font=Font.HEAD, size=34, bold=True, color=Color.FOREGROUND,
        letter_spacing=-0.01,
    )
    add_text(
        s, MARGIN_L, 2.65, 7.5, 0.7,
        promise,
        font=Font.BODY, size=16, color=Color.SECONDARY_FG,
    )

    add_bulleted_list(
        s, MARGIN_L, 3.55, 7.3, 3.0,
        bullets, font=Font.BODY, size=14, color=Color.SECONDARY_FG,
        bullet_color=Color.GOLD, bullet_gap=0.62,
    )

    # Right column visual card
    card_x, card_y, card_w, card_h = 8.55, 1.35, 4.2, 5.05
    add_card(s, card_x, card_y, card_w, card_h)
    _render_visual(s, visual_kind, visual_payload or {}, card_x, card_y, card_w, card_h)


def _render_visual(s, kind: str, payload: dict, cx: float, cy: float, cw: float, ch: float):
    """Render the right-column visual within a card."""
    inner_x = cx + 0.3
    inner_y = cy + 0.4
    inner_w = cw - 0.6
    inner_h = ch - 0.8

    # Card overline (top-left of every visual)
    add_text(
        s, inner_x, inner_y - 0.05, inner_w, 0.3,
        payload.get("label", "PREVIEW"),
        font=Font.HEAD, size=9, bold=True, color=Color.MUTED,
        letter_spacing=0.15,
    )

    if kind == "brand_stack":
        # 4 stacked brand tiles
        names = payload.get("names", ["BRAND ALPHA", "BRAND BETA", "BRAND GAMMA", "BRAND DELTA"])
        tile_h = 0.62
        gap = 0.14
        top = inner_y + 0.5
        for i, name in enumerate(names):
            ty = top + i * (tile_h + gap)
            add_rect(s, inner_x, ty, inner_w, tile_h, fill=Color.MUTED_BG, line=GoldOn.CARD_12)
            # Color dot to suggest brand
            dot_colors = [Color.GOLD, Color.CYAN, Color.CORAL, Color.LIME]
            add_rect(s, inner_x + 0.18, ty + 0.22, 0.18, 0.18, fill=dot_colors[i % 4])
            add_text(
                s, inner_x + 0.5, ty + 0.13, inner_w - 0.6, 0.4,
                name, font=Font.HEAD, size=11, bold=True, color=Color.FOREGROUND,
                letter_spacing=0.10,
            )
            add_text(
                s, inner_x + 0.5, ty + 0.34, inner_w - 0.6, 0.3,
                f"{(i + 1) * 7}.{(i * 13) % 10}M  followers",
                font=Font.MONO, size=8, color=Color.MUTED, letter_spacing=0.04,
            )
        return

    if kind == "logo_grid":
        # 3 x 2 logo tiles
        cols, rows = 3, 2
        tile_w = (inner_w - 0.2 * (cols - 1)) / cols
        tile_h = (inner_h - 0.4 - 0.2 * (rows - 1)) / rows
        labels = payload.get("labels", ["A", "B", "C", "D", "E", "F"])
        for r in range(rows):
            for c in range(cols):
                x = inner_x + c * (tile_w + 0.2)
                y = inner_y + 0.5 + r * (tile_h + 0.2)
                add_rect(s, x, y, tile_w, tile_h, fill=Color.MUTED_BG, line=GoldOn.CARD_12)
                add_text(
                    s, x, y, tile_w, tile_h,
                    labels[r * cols + c],
                    font=Font.DISPLAY, size=36, bold=True, color=Color.GOLD,
                    align="center", anchor="middle",
                )
        return

    if kind == "product_card":
        # A product card mock — square image + title + price
        pad = 0.3
        img_h = inner_h * 0.55
        add_rect(s, inner_x, inner_y + 0.5, inner_w, img_h, fill=Color.MUTED_BG, line=GoldOn.CARD_12)
        # Diagonal corner to suggest content
        add_rect(s, inner_x + 0.4, inner_y + 0.5 + img_h - 1.2, inner_w - 0.8, 0.04, fill=Color.GOLD)
        add_text(
            s, inner_x, inner_y + 0.5 + img_h - 1.0, inner_w, 0.3,
            "RUNNER · LIMITED",
            font=Font.HEAD, size=10, bold=True, color=Color.GOLD,
            align="center", letter_spacing=0.20,
        )
        add_text(
            s, inner_x, inner_y + 0.5 + img_h + 0.2, inner_w, 0.5,
            payload.get("title", "Stratus Lo Sneaker"),
            font=Font.HEAD, size=18, bold=True, color=Color.FOREGROUND, align="center",
        )
        add_text(
            s, inner_x, inner_y + 0.5 + img_h + 0.75, inner_w, 0.35,
            payload.get("subtitle", "Performance · SS26"),
            font=Font.BODY, size=11, color=Color.MUTED, align="center",
        )
        add_text(
            s, inner_x, inner_y + 0.5 + img_h + 1.15, inner_w, 0.4,
            payload.get("price", "$ 189.00"),
            font=Font.MONO, size=18, bold=True, color=Color.GOLD,
            align="center", letter_spacing=0.04,
        )
        return

    if kind == "template_diagram":
        # A template layout mock — three zones
        add_rect(s, inner_x, inner_y + 0.5, inner_w, 1.6, fill=Color.MUTED_BG, line=GoldOn.CARD_12)
        add_text(
            s, inner_x, inner_y + 0.5, inner_w, 1.6,
            "PRODUCT HERO",
            font=Font.HEAD, size=11, bold=True, color=Color.MUTED,
            align="center", anchor="middle", letter_spacing=0.20,
        )
        add_rect(s, inner_x, inner_y + 2.25, inner_w * 0.6 - 0.05, 0.55, fill=Color.MUTED_BG, line=GoldOn.CARD_12)
        add_text(
            s, inner_x, inner_y + 2.25, inner_w * 0.6 - 0.05, 0.55,
            "{{TITLE}}", font=Font.MONO, size=10, color=Color.GOLD,
            align="center", anchor="middle",
        )
        add_rect(s, inner_x + inner_w * 0.6 + 0.05, inner_y + 2.25, inner_w * 0.4 - 0.05, 0.55,
                 fill=Color.MUTED_BG, line=GoldOn.CARD_12)
        add_text(
            s, inner_x + inner_w * 0.6 + 0.05, inner_y + 2.25, inner_w * 0.4 - 0.05, 0.55,
            "{{PRICE}}", font=Font.MONO, size=10, color=Color.GOLD,
            align="center", anchor="middle",
        )
        add_rect(s, inner_x, inner_y + 2.95, inner_w, 0.55, fill=Color.MUTED_BG, line=GoldOn.CARD_12)
        add_text(
            s, inner_x, inner_y + 2.95, inner_w, 0.55,
            "{{CTA}}", font=Font.MONO, size=10, color=Color.GOLD,
            align="center", anchor="middle",
        )
        # Caption
        add_text(
            s, inner_x, inner_y + 3.65, inner_w, 0.3,
            "FIELDS LOCK BRAND CONSISTENCY",
            font=Font.HEAD, size=8, bold=True, color=Color.MUTED,
            align="center", letter_spacing=0.15,
        )
        return

    if kind == "swatch_stack":
        # 4 horizontal color swatches — chart palette
        colors = [Color.GOLD, Color.CYAN, Color.CORAL, Color.LIME]
        names = ["GOLD", "CYAN", "CORAL", "LIME"]
        hexes = ["#F4B964", "#64DCF4", "#E8956A", "#A4F464"]
        swatch_h = 0.62
        gap = 0.14
        top = inner_y + 0.5
        for i, c in enumerate(colors):
            ty = top + i * (swatch_h + gap)
            add_rect(s, inner_x, ty, 0.7, swatch_h, fill=c)
            add_rect(s, inner_x + 0.7, ty, inner_w - 0.7, swatch_h,
                     fill=Color.MUTED_BG, line=GoldOn.CARD_12)
            add_text(
                s, inner_x + 0.9, ty + 0.05, inner_w - 1.0, 0.32,
                names[i], font=Font.HEAD, size=10, bold=True, color=Color.FOREGROUND,
                letter_spacing=0.15,
            )
            add_text(
                s, inner_x + 0.9, ty + 0.3, inner_w - 1.0, 0.32,
                hexes[i], font=Font.MONO, size=9, color=Color.MUTED, letter_spacing=0.04,
            )
        return

    if kind == "caption_mock":
        # A caption text block with line indicators
        add_text(
            s, inner_x, inner_y + 0.55, inner_w, 0.4,
            "DRAFT 03 / TONE: WARM",
            font=Font.HEAD, size=9, bold=True, color=Color.GOLD,
            letter_spacing=0.18,
        )
        add_paragraphs(
            s, inner_x, inner_y + 1.0, inner_w, 2.6,
            [
                "Made for the long route home.",
                "Engineered to shrug off everything",
                "the sidewalk has in mind.",
                "",
                "#stratus  #ss26  #everydaycarry",
            ],
            font=Font.BODY, size=13, color=Color.FOREGROUND, line_spacing=1.5,
        )
        # Tone meter
        add_text(
            s, inner_x, inner_y + 3.65, inner_w, 0.3,
            "TONE MATCH",
            font=Font.HEAD, size=8, bold=True, color=Color.MUTED,
            letter_spacing=0.15,
        )
        add_rect(s, inner_x, inner_y + 3.9, inner_w, 0.08, fill=Color.MUTED_BG)
        add_rect(s, inner_x, inner_y + 3.9, inner_w * 0.86, 0.08, fill=Color.GOLD)
        return

    if kind == "approval_flow":
        # Three approval states in a vertical stack
        states = [
            ("DRAFT",     "Created by Maya",     Color.MUTED),
            ("REVIEW",    "Awaiting Brand Lead", Color.CORAL),
            ("APPROVED",  "Ready to schedule",   Color.LIME),
        ]
        tile_h = 0.8
        gap = 0.18
        top = inner_y + 0.55
        for i, (state, sub, dot_c) in enumerate(states):
            ty = top + i * (tile_h + gap)
            add_rect(s, inner_x, ty, inner_w, tile_h, fill=Color.MUTED_BG, line=GoldOn.CARD_12)
            add_rect(s, inner_x + 0.2, ty + 0.32, 0.18, 0.18, fill=dot_c)
            add_text(
                s, inner_x + 0.55, ty + 0.16, inner_w - 0.7, 0.35,
                state, font=Font.HEAD, size=12, bold=True, color=Color.FOREGROUND,
                letter_spacing=0.12,
            )
            add_text(
                s, inner_x + 0.55, ty + 0.43, inner_w - 0.7, 0.3,
                sub, font=Font.BODY, size=10, color=Color.MUTED,
            )
        return

    if kind == "calendar_mock":
        # Mini calendar grid (5 cols x 3 rows)
        cols, rows = 5, 3
        labels = ["MON", "TUE", "WED", "THU", "FRI"]
        for c, lbl in enumerate(labels):
            x = inner_x + c * (inner_w / cols)
            add_text(
                s, x, inner_y + 0.5, inner_w / cols, 0.3,
                lbl, font=Font.HEAD, size=9, bold=True, color=Color.MUTED,
                align="center", letter_spacing=0.18,
            )
        tile_w = (inner_w / cols) - 0.08
        tile_h = (inner_h - 1.2) / rows - 0.1
        # Mark certain tiles as "scheduled"
        scheduled = {(0, 1), (1, 3), (2, 0), (2, 2), (2, 4), (0, 4)}
        for r in range(rows):
            for c in range(cols):
                x = inner_x + c * (inner_w / cols) + 0.04
                y = inner_y + 0.85 + r * (tile_h + 0.1)
                is_sched = (r, c) in scheduled
                add_rect(
                    s, x, y, tile_w, tile_h,
                    fill=Color.MUTED_BG if not is_sched else GoldOn.CARD_22,
                    line=GoldOn.CARD_12,
                )
                if is_sched:
                    add_rect(s, x + 0.1, y + 0.1, 0.12, 0.12, fill=Color.GOLD)
        # Footer note
        add_text(
            s, inner_x, inner_y + 3.55, inner_w, 0.3,
            "06 POSTS QUEUED · NEXT WEEK",
            font=Font.HEAD, size=8, bold=True, color=Color.MUTED,
            align="center", letter_spacing=0.18,
        )
        return

    if kind == "platforms":
        platforms = [
            ("INSTAGRAM",  "LIVE",      Color.GOLD),
            ("TIKTOK",     "ROADMAP",   Color.MUTED),
            ("LINKEDIN",   "ROADMAP",   Color.MUTED),
            ("X",          "ROADMAP",   Color.MUTED),
        ]
        tile_h = 0.74
        gap = 0.16
        top = inner_y + 0.55
        for i, (name, badge, badge_c) in enumerate(platforms):
            ty = top + i * (tile_h + gap)
            add_rect(s, inner_x, ty, inner_w, tile_h, fill=Color.MUTED_BG, line=GoldOn.CARD_12)
            add_text(
                s, inner_x + 0.25, ty + 0.18, inner_w - 0.5, 0.45,
                name, font=Font.HEAD, size=14, bold=True, color=Color.FOREGROUND,
                letter_spacing=0.10,
            )
            # Badge
            badge_w = 0.85
            badge_x = inner_x + inner_w - badge_w - 0.2
            add_rect(s, badge_x, ty + 0.22, badge_w, 0.32,
                     fill=Color.MUTED_BG, line=badge_c)
            add_text(
                s, badge_x, ty + 0.22, badge_w, 0.32,
                badge, font=Font.HEAD, size=8, bold=True, color=badge_c,
                align="center", anchor="middle", letter_spacing=0.18,
            )
        return

    if kind == "stat_panel":
        # Top: a big mono number
        big = payload.get("stat", "+34%")
        cap = payload.get("stat_caption", "median engagement lift")
        add_text(
            s, inner_x, inner_y + 0.55, inner_w, 1.4,
            big, font=Font.MONO, size=78, bold=True, color=Color.GOLD,
            align="left",
        )
        add_text(
            s, inner_x, inner_y + 2.0, inner_w, 0.4,
            cap, font=Font.BODY, size=13, color=Color.SECONDARY_FG,
        )
        add_gold_rule(s, inner_x, inner_y + 2.55, 1.2)
        # Sub-stats
        sub = payload.get("substats", [
            ("FOLLOWER GROWTH",   "+ 12.4%"),
            ("REACH",             "+ 41.0%"),
            ("SAVES PER POST",    "+ 2.3 ×"),
        ])
        for i, (label, val) in enumerate(sub):
            y = inner_y + 2.75 + i * 0.42
            add_text(
                s, inner_x, y, inner_w * 0.6, 0.35,
                label, font=Font.HEAD, size=9, bold=True, color=Color.MUTED,
                letter_spacing=0.15,
            )
            add_text(
                s, inner_x + inner_w * 0.55, y, inner_w * 0.45, 0.35,
                val, font=Font.MONO, size=12, bold=True, color=Color.FOREGROUND,
                align="right",
            )
        return

    if kind == "enterprise_grid":
        items = payload.get("items", [
            ("ISOLATED WORKSPACES",  "Per-brand boundaries"),
            ("ROLE-BASED ACCESS",    "Granular permissions"),
            ("SECURE AUTH",          "SOC-aligned via Clerk"),
            ("AUDIT TRAIL",          "Every action logged"),
        ])
        cols, rows = 2, 2
        tile_w = (inner_w - 0.2) / cols
        tile_h = (inner_h - 0.5 - 0.2) / rows
        for i, (title, sub) in enumerate(items):
            r, c = divmod(i, cols)
            x = inner_x + c * (tile_w + 0.2)
            y = inner_y + 0.5 + r * (tile_h + 0.2)
            add_rect(s, x, y, tile_w, tile_h, fill=Color.MUTED_BG, line=GoldOn.CARD_12)
            add_rect(s, x, y, 0.06, tile_h, fill=Color.GOLD)
            add_text(
                s, x + 0.2, y + 0.2, tile_w - 0.3, 0.4,
                title, font=Font.HEAD, size=11, bold=True, color=Color.FOREGROUND,
                letter_spacing=0.12,
            )
            add_text(
                s, x + 0.2, y + 0.55, tile_w - 0.3, 0.4,
                sub, font=Font.BODY, size=10, color=Color.MUTED,
            )


def build_30_day_plan(prs, page_num):
    s = slide_frame(prs, page_num)
    add_overline(s, MARGIN_L, 0.7, 8, "04 / Onboarding Plan")
    add_gold_accent_bar(s, MARGIN_L, 1.4, height=0.65)
    add_text(
        s, MARGIN_L + 0.25, 1.35, 12, 1.0,
        "Your first 30 days with PixelPrism.",
        font=Font.HEAD, size=34, bold=True, color=Color.FOREGROUND,
        letter_spacing=-0.01,
    )
    add_text(
        s, MARGIN_L, 2.65, 11.5, 0.5,
        "Four weeks from kickoff to compounding output.",
        font=Font.BODY, size=15, color=Color.SECONDARY_FG,
    )

    weeks = [
        ("WEEK 01", "CONFIGURE", "Spin up workspaces. Upload logos, products, brand voice. Assign roles."),
        ("WEEK 02", "CREATE",    "First AI-generated campaigns. Train tone of voice. Approve early drafts."),
        ("WEEK 03", "CODIFY",    "Lock in templates and approval flow. Document your team's playbook."),
        ("WEEK 04", "COMPOUND",  "Full publishing cadence live. First analytics review. Plan next sprint."),
    ]
    card_y = 3.7
    card_h = 2.9
    gap = 0.22
    card_w = (SLIDE_W - 2 * MARGIN_L - 3 * gap) / 4

    # Connecting gold line through the cards
    line_y = card_y + 0.55
    add_rect(s, MARGIN_L + 0.25, line_y, SLIDE_W - 2 * MARGIN_L - 0.5, 0.025, fill=GoldOn.BG_22)

    for i, (week, verb, desc) in enumerate(weeks):
        x = MARGIN_L + i * (card_w + gap)
        # Node on the line
        node_x = x + card_w / 2 - 0.07
        add_rect(s, node_x, line_y - 0.05, 0.14, 0.14, fill=Color.GOLD)
        # Card
        add_card(s, x, card_y + 0.8, card_w, card_h - 0.8)
        # Week label
        add_text(
            s, x, card_y, card_w, 0.4,
            week, font=Font.MONO, size=10, bold=True, color=Color.MUTED,
            align="center", letter_spacing=0.18,
        )
        # Verb (big)
        add_text(
            s, x + 0.2, card_y + 1.0, card_w - 0.4, 0.7,
            verb, font=Font.HEAD, size=22, bold=True, color=Color.GOLD,
            letter_spacing=0.05,
        )
        # Description
        add_paragraphs(
            s, x + 0.2, card_y + 1.75, card_w - 0.4, 1.4,
            [desc],
            font=Font.BODY, size=11, color=Color.SECONDARY_FG, line_spacing=1.5,
        )


def build_support(prs, page_num):
    s = slide_frame(prs, page_num)
    add_overline(s, MARGIN_L, 0.7, 8, "05 / Support & Contact")
    add_gold_accent_bar(s, MARGIN_L, 1.4, height=0.65)
    add_text(
        s, MARGIN_L + 0.25, 1.35, 12, 1.0,
        "We're on the call when you need us.",
        font=Font.HEAD, size=34, bold=True, color=Color.FOREGROUND,
        letter_spacing=-0.01,
    )
    add_text(
        s, MARGIN_L, 2.65, 7.5, 0.5,
        "Your kickoff is the start of an ongoing partnership.",
        font=Font.BODY, size=15, color=Color.SECONDARY_FG,
    )

    add_bulleted_list(
        s, MARGIN_L, 3.5, 7.3, 3.0,
        [
            "Dedicated kickoff lead for the first 30 days.",
            "Shared Slack channel with PixelPrism engineering.",
            "Weekly office hours. Escalation acknowledged within 4 hours.",
            "Roadmap input — your team helps prioritize what we build next.",
        ],
        font=Font.BODY, size=14, color=Color.SECONDARY_FG, bullet_gap=0.62,
    )

    # Contact card on the right
    add_card(s, 8.55, 2.8, 4.2, 3.6)
    add_text(
        s, 8.85, 3.0, 3.6, 0.35,
        "REACH US",
        font=Font.HEAD, size=10, bold=True, color=Color.GOLD,
        letter_spacing=0.18,
    )
    contacts = [
        ("KICKOFF LEAD", "kickoff@pixelprism.app"),
        ("ENGINEERING",  "eng@pixelprism.app"),
        ("DOCS",         "pixelprism.app/docs"),
        ("STATUS",       "pixelprism.app/status"),
    ]
    for i, (label, val) in enumerate(contacts):
        y = 3.5 + i * 0.7
        add_text(
            s, 8.85, y, 3.6, 0.3,
            label, font=Font.HEAD, size=9, bold=True, color=Color.MUTED,
            letter_spacing=0.18,
        )
        add_text(
            s, 8.85, y + 0.3, 3.6, 0.35,
            val, font=Font.MONO, size=12, bold=True, color=Color.FOREGROUND,
            letter_spacing=0.03,
        )


def build_closing(prs, page_num):
    s = slide_frame(prs, page_num)
    add_overline(s, MARGIN_L, 0.55, 8, "End / Thank You")
    # Big Display closer
    add_text(
        s, MARGIN_L, 2.1, 12, 1.4,
        "Let's build",
        font=Font.DISPLAY, size=82, bold=True, color=Color.SECONDARY_FG,
        letter_spacing=-0.02,
    )
    add_text(
        s, MARGIN_L, 3.3, 12, 1.4,
        "something",
        font=Font.DISPLAY, size=82, bold=True, color=Color.FOREGROUND,
        letter_spacing=-0.02,
    )
    add_text(
        s, MARGIN_L, 4.5, 12, 1.4,
        "sharp.",
        font=Font.DISPLAY, size=82, bold=True, color=Color.GOLD,
        letter_spacing=-0.02,
    )
    add_gold_rule(s, MARGIN_L, 6.05, 2.4)
    add_text(
        s, MARGIN_L, 6.25, 6, 0.35,
        "PIXELPRISM · ENTERPRISE KICKOFF",
        font=Font.HEAD, size=10, bold=True, color=Color.MUTED,
        letter_spacing=0.20,
    )
    add_text(
        s, SLIDE_W - 4, 6.25, 3.4, 0.35,
        "v1.0 · 2026",
        font=Font.MONO, size=10, color=Color.MUTED,
        align="right", letter_spacing=0.04,
    )


# ============================================================================
# DECK ASSEMBLY
# ============================================================================

def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # 01
    build_cover(prs, 1)
    # 02
    build_problem(prs, 2)
    # 03
    build_what_is_pp(prs, 3)
    # 04 — Act I divider
    build_act_divider(
        prs, 4, "I", "Set Up Your Brand.",
        "Configure once. Reuse forever.",
    )
    # 05
    build_feature(
        prs, 5,
        overline="Act I / Set Up",
        title="One workspace, many brands.",
        promise="Each brand is its own fully-isolated world inside PixelPrism.",
        bullets=[
            "Spin up dedicated workspaces per brand — separate assets, voice, and team.",
            "Switch contexts in one click. No cross-contamination between brands.",
            "Role-based access per brand. Your apparel team never sees your beauty team's work.",
        ],
        visual_kind="brand_stack",
        visual_payload={
            "label": "BRAND WORKSPACES",
            "names": ["STRATUS APPAREL", "AURA BEAUTY", "NORTH WORKS", "FERN & OAK"],
        },
    )
    # 06
    build_feature(
        prs, 6,
        overline="Act I / Set Up",
        title="Your brand's library, always at hand.",
        promise="Upload your identity once. The system uses it everywhere you create.",
        bullets=[
            "Logos, color palettes, and reference imagery live in one place.",
            "Every generated asset pulls from the library automatically.",
            "No more “where's the latest logo?” Slack threads at 4pm.",
        ],
        visual_kind="logo_grid",
        visual_payload={"label": "BRAND LIBRARY", "labels": ["S", "A", "N", "F", "+", "+"]},
    )
    # 07
    build_feature(
        prs, 7,
        overline="Act I / Set Up",
        title="Your catalog, weaponized for content.",
        promise="Every product becomes a first-class content primitive.",
        bullets=[
            "Each product has an image, story, and metadata in PixelPrism.",
            "Drag a product into a template — fields auto-fill. No copy-paste.",
            "When the catalog updates, content updates with it.",
        ],
        visual_kind="product_card",
        visual_payload={
            "label": "PRODUCT RECORD",
            "title": "Stratus Lo Sneaker",
            "subtitle": "Performance · SS26",
            "price": "$ 189.00",
        },
    )
    # 08
    build_feature(
        prs, 8,
        overline="Act I / Set Up",
        title="Templates that lock in your brand.",
        promise="Brand-fielded layouts so no one can drift off-brand by accident.",
        bullets=[
            "Define your most-used post formats — hero shot, product feature, quote card, carousel intro.",
            "Each template is a layout with brand-locked fields you fill in.",
            "Onboard a new designer in an hour. The template won't let them break the brand.",
        ],
        visual_kind="template_diagram",
        visual_payload={"label": "TEMPLATE SCHEMA"},
    )
    # 09 — Act II divider
    build_act_divider(
        prs, 9, "II", "Make the Work.",
        "Generate. Refine. Approve.",
    )
    # 10
    build_feature(
        prs, 10,
        overline="Act II / Create",
        title="Studio-grade visuals, generated to brief.",
        promise="AI image generation tuned to your brand colors, logos, and product imagery.",
        bullets=[
            "Prompt with intent. (“Hero shot, new runner, wet pavement, dusk.”)",
            "Output respects your palette, your logos, your product shapes.",
            "Iterate in seconds. Save the keepers. Discard the rest.",
        ],
        visual_kind="swatch_stack",
        visual_payload={"label": "BRAND PALETTE"},
    )
    # 11
    build_feature(
        prs, 11,
        overline="Act II / Create",
        title="Captions that sound like your brand.",
        promise="A caption generator that learns your tone — and avoids the tones you've rejected.",
        bullets=[
            "Generate captions, hooks, and hashtag sets in your voice.",
            "Trained on what you've approved. Avoids patterns you've rejected.",
            "Built into the scheduler — no more shuffling between tools.",
        ],
        visual_kind="caption_mock",
        visual_payload={"label": "CAPTION DRAFT"},
    )
    # 12
    build_feature(
        prs, 12,
        overline="Act II / Create",
        title="Approval, without the email chain.",
        promise="A real review queue per brand — with a clear audit trail.",
        bullets=[
            "Drafts queue per brand. Reviewers approve or comment inline.",
            "Approved work moves to scheduling automatically.",
            "Full audit trail: who said yes, who said no, when.",
        ],
        visual_kind="approval_flow",
        visual_payload={"label": "REVIEW QUEUE"},
    )
    # 13 — Act III divider
    build_act_divider(
        prs, 13, "III", "Publish & Learn.",
        "Ship the work. Measure the lift.",
    )
    # 14
    build_feature(
        prs, 14,
        overline="Act III / Publish",
        title="Your Instagram calendar, in one place.",
        promise="Plan the week. Drop in posts. Publish without leaving the studio.",
        bullets=[
            "Drag-drop scheduling per brand. See gaps before they happen.",
            "Native Instagram publishing — Feed, Reels, and Stories from one composer.",
            "Reschedule, swap captions, or pause campaigns in a click.",
        ],
        visual_kind="calendar_mock",
        visual_payload={"label": "WEEKLY SCHEDULE"},
    )
    # 15
    build_feature(
        prs, 15,
        overline="Act III / Publish",
        title="Instagram first. Built to expand.",
        promise="Today: Instagram. Tomorrow: every channel your team needs.",
        bullets=[
            "Instagram is the launch target. The scheduler is platform-agnostic underneath.",
            "TikTok, LinkedIn, and X integrations live on the roadmap.",
            "Same composer, same approvals — more channels, no extra workflow.",
        ],
        visual_kind="platforms",
        visual_payload={"label": "CHANNELS"},
    )
    # 16
    build_feature(
        prs, 16,
        overline="Act III / Learn",
        title="Performance, per brand.",
        promise="Engagement, reach, and growth — sliced by brand, post type, and campaign.",
        bullets=[
            "See what's working before the week ends, not after the quarter.",
            "Compare brands, periods, and post formats in one view.",
            "Export raw metrics to whatever BI stack you already run.",
        ],
        visual_kind="stat_panel",
        visual_payload={
            "label": "PILOT RESULTS",
            "stat": "+34%",
            "stat_caption": "median engagement lift across pilot brands",
            "substats": [
                ("FOLLOWER GROWTH",   "+ 12.4 %"),
                ("REACH",             "+ 41.0 %"),
                ("SAVES PER POST",    "+ 2.3 ×"),
            ],
        },
    )
    # 17
    build_feature(
        prs, 17,
        overline="03 / Enterprise Readiness",
        title="Built for teams that move serious volume.",
        promise="Multi-brand isolation, role-based access, and audit trails out of the box.",
        bullets=[
            "Strict per-brand workspace isolation. Nothing leaks across brands.",
            "Role-based permissions tied to SSO-grade auth (Clerk).",
            "Audit trails on creative, approvals, and publishing.",
            "Per-seat or per-brand pricing. Let's tailor it.",
        ],
        visual_kind="enterprise_grid",
        visual_payload={
            "label": "ENTERPRISE CONTROLS",
            "items": [
                ("ISOLATED WORKSPACES", "Per-brand boundaries"),
                ("ROLE-BASED ACCESS",   "Granular permissions"),
                ("SECURE AUTH",         "SSO-grade via Clerk"),
                ("AUDIT TRAIL",         "Every action logged"),
            ],
        },
    )
    # 18
    build_30_day_plan(prs, 18)
    # 19
    build_support(prs, 19)
    # 20
    build_closing(prs, 20)

    out_path = os.path.join(os.path.dirname(__file__), "PixelPrism-Enterprise-Kickoff.pptx")
    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    path = build_deck()
    print(f"Wrote {path}")
